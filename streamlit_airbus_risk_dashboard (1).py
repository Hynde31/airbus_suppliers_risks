Streamlit Airbus Risk Dashboard
Introduction
Cette application Streamlit est conçue pour la gestion des risques fournisseurs dans la supply chain aéronautique, en particulier pour Airbus. Elle intègre des fonctionnalités avancées pour simuler des scénarios géopolitiques, évaluer dynamiquement les risques, prédire les ruptures de stock, cartographier les dépendances critiques, et générer des rapports PowerPoint.
Code

import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime, timedelta

# Sample data for suppliers
data = {
    'Supplier': ['Supplier A', 'Supplier B', 'Supplier C'],
    'Country': ['France', 'Germany', 'USA'],
    'Risk Score': [0.2, 0.5, 0.8],
    'Risk Level': ['Low', 'Medium', 'High'],
    'Average Delay': [5, 10, 15],
    'Stock Level': [100, 50, 20],
    'Incidents': [1, 3, 5],
    'latitude': [48.8566, 52.5200, 37.7749],
    'longitude': [2.3522, 13.4050, -122.4194]
}

# Convert data to DataFrame
df = pd.DataFrame(data)

# Streamlit app
st.title('Supplier Risk Management Platform')

# Map visualization
st.map(df)

# Supplier selector
supplier = st.selectbox('Select a supplier:', df['Supplier'])

# Display supplier details
supplier_data = df[df['Supplier'] == supplier].iloc[0]
st.write(f"**Country:** {supplier_data['Country']}")
st.write(f"**Risk Score:** {supplier_data['Risk Score']}")
st.write(f"**Risk Level:** {supplier_data['Risk Level']}")
st.write(f"**Average Delay:** {supplier_data['Average Delay']} days")
st.write(f"**Stock Level:** {supplier_data['Stock Level']}")
st.write(f"**Incidents:** {supplier_data['Incidents']}")

# Dashboard
st.bar_chart(df[['Supplier', 'Risk Score']].set_index('Supplier'))
st.line_chart(df[['Supplier', 'Average Delay']].set_index('Supplier'))
st.area_chart(df[['Supplier', 'Stock Level']].set_index('Supplier'))

# Simulation de scénarios géopolitiques
st.sidebar.title('Simulation de scénarios géopolitiques')
scenario = st.sidebar.selectbox('Sélectionnez un scénario:', ['Aucun', 'Embargo', 'Grève', 'Guerre'])

if scenario == 'Embargo':
    df['Risk Score'] += 0.1
    df['Average Delay'] += 5
elif scenario == 'Grève':
    df['Risk Score'] += 0.2
    df['Average Delay'] += 10
elif scenario == 'Guerre':
    df['Risk Score'] += 0.3
    df['Average Delay'] += 15

# Scoring dynamique avec actualités (simulé)
impact = st.sidebar.slider('Impact des actualités sur le risque:', 0.0, 1.0, 0.1)
df['Risk Score'] += impact

# Analyse prédictive des ruptures de stock
df['Predicted Stock Out'] = df['Average Delay'] * 2

# Cartographie des dépendances critiques (BOM)
critical_components = {
    'Supplier A': ['Moteur', 'Avionique'],
    'Supplier B': ['Fuselage', 'Train d'atterrissage'],
    'Supplier C': ['Systèmes électriques', 'Hydraulique']
}

st.sidebar.title('Composants critiques')
selected_supplier = st.sidebar.selectbox('Sélectionnez un fournisseur:', df['Supplier'])
st.sidebar.write(f"Composants critiques pour {selected_supplier}: {', '.join(critical_components[selected_supplier])}")

# Dashboard par programme avion
program_data = {
    'Programme': ['A320', 'A350', 'A220'],
    'Risk Score': [0.4, 0.6, 0.3],
    'Average Delay': [7, 12, 5],
    'Stock Level': [80, 40, 60]
}

df_program = pd.DataFrame(program_data)
st.title('Dashboard par programme avion')
st.bar_chart(df_program[['Programme', 'Risk Score']].set_index('Programme'))
st.line_chart(df_program[['Programme', 'Average Delay']].set_index('Programme'))
st.area_chart(df_program[['Programme', 'Stock Level']].set_index('Programme'))

# Export PowerPoint
def export_ppt():
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Rapport de gestion des risques fournisseurs"

    # Logo Airbus
    slide.shapes.add_picture('airbus_logo.png', Inches(0.5), Inches(0.5), height=Inches(1.0))

    # Détails du fournisseur
    slide.placeholders[1].text = (
        f"Fournisseur: {supplier_data['Supplier']}
"
        f"Pays: {supplier_data['Country']}
"
        f"Score de risque: {supplier_data['Risk Score']}
"
        f"Niveau de risque: {supplier_data['Risk Level']}
"
        f"Délai moyen: {supplier_data['Average Delay']} jours
"
        f"Niveau de stock: {supplier_data['Stock Level']}
"
        f"Incidents: {supplier_data['Incidents']}
"
    )

    # Recommandations
    slide.placeholders[1].text += "
Recommandations:
"
    if supplier_data['Risk Level'] == 'High':
        slide.placeholders[1].text += "- Diversifier les fournisseurs
"
        slide.placeholders[1].text += "- Augmenter les stocks de sécurité
"
    elif supplier_data['Risk Level'] == 'Medium':
        slide.placeholders[1].text += "- Surveiller les performances
"
        slide.placeholders[1].text += "- Planifier des audits réguliers
"
    else:
        slide.placeholders[1].text += "- Maintenir la collaboration
"

    prs.save('rapport_risques_fournisseurs.pptx')
    st.success('Rapport PowerPoint généré avec succès!')

if st.button('Exporter en PowerPoint'):
    export_ppt()

