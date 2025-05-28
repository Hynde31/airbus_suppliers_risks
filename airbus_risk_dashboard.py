import streamlit as st
import pandas as pd
import numpy as np
import os
from datetime import datetime

# Gestion import pptx (PowerPoint)
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from suppliers_data import SUPPLIERS

st.set_page_config(page_title="Airbus Suppliers Risk Dashboard", layout="wide")

# --------- DATA FLATTENING ---------
def flatten_suppliers(suppliers):
    rows = []
    for s in suppliers:
        for site in s["sites"]:
            rows.append({
                "Supplier": s["name"],
                "Component": s["component"],
                "Criticality": s["criticality"],
                "Dual Sourcing": "Yes" if s["dual_sourcing"] else "No",
                "Site City": site["city"],
                "Country": site["country"],
                "latitude": site["lat"],
                "longitude": site["lon"],
                "Stock Days": site["stock_days"],
                "Lead Time": site["lead_time"],
                "On-Time Delivery (%)": site.get("on_time_delivery", 90),
                "Incidents": site.get("incidents", 1)
            })
    return pd.DataFrame(rows)

df = flatten_suppliers(SUPPLIERS)

# --- KPI DERIVED COLUMNS ---
df["Risk Score"] = (1 - df["On-Time Delivery (%)"]/100) \
                   + (df["Incidents"]/10) \
                   + (1 - df["Stock Days"]/60)*0.5 \
                   + (df["Lead Time"]/150)*0.5
df["Risk Level"] = pd.cut(df["Risk Score"], bins=[-np.inf,0.5,1,2], labels=["Low","Medium","High"])

st.title("Airbus Suppliers Risk Dashboard")
logo_path = "airbus_logo.png"
if os.path.exists(logo_path):
    st.image(logo_path, width=180)

# --------- SIDEBAR: SCENARIO & FILTERS ---------
st.sidebar.header("Scenario Simulation & Filters")
scenario = st.sidebar.selectbox(
    'Geopolitical Scenario Simulation',
    ['None', 'Embargo', 'Strike', 'War']
)
impact = st.sidebar.slider('News impact on risk (adds to "Incidents"):', 0, 5, 0)
selected_criticality = st.sidebar.multiselect(
    "Criticality", options=sorted(df["Criticality"].unique()), default=list(df["Criticality"].unique()))
selected_country = st.sidebar.multiselect(
    "Country", options=sorted(df["Country"].unique()), default=list(df["Country"].unique()))
selected_component = st.sidebar.multiselect(
    "Component", options=sorted(df["Component"].unique()), default=list(df["Component"].unique()))

# --------- SCENARIO IMPACT ---------
df_scenar = df.copy()
if scenario == "Embargo":
    df_scenar["Lead Time"] += 15
    df_scenar["Incidents"] += 2
elif scenario == "Strike":
    df_scenar["Lead Time"] += 10
    df_scenar["Incidents"] += 3
elif scenario == "War":
    df_scenar["Lead Time"] += 30
    df_scenar["Incidents"] += 6
if impact:
    df_scenar["Incidents"] += impact

df_scenar["Risk Score"] = (1 - df_scenar["On-Time Delivery (%)"]/100) \
                          + (df_scenar["Incidents"]/10) \
                          + (1 - df_scenar["Stock Days"]/60)*0.5 \
                          + (df_scenar["Lead Time"]/150)*0.5
df_scenar["Risk Level"] = pd.cut(df_scenar["Risk Score"], bins=[-np.inf,0.5,1,2], labels=["Low","Medium","High"])

# --------- FILTRAGE ---------
filtered_df = df_scenar[
    df_scenar["Criticality"].isin(selected_criticality)
    & df_scenar["Country"].isin(selected_country)
    & df_scenar["Component"].isin(selected_component)
]

# --------- MAP OF SUPPLIER SITES ---------
st.subheader("Supplier Sites Map")
if (
    not filtered_df.empty and 
    "latitude" in filtered_df.columns and 
    "longitude" in filtered_df.columns and
    filtered_df[["latitude", "longitude"]].notnull().all().all()
):
    st.map(filtered_df[["latitude", "longitude"]])
else:
    st.info("No supplier sites match your filters or missing coordinates.")

# --------- SUPPLIERS TABLE ---------
st.subheader("Suppliers Table")
st.dataframe(filtered_df, use_container_width=True)

# --------- GLOBAL STATISTICS ---------
st.subheader("Global Statistics")
col1, col2, col3 = st.columns(3)
col1.metric("Unique Suppliers", filtered_df["Supplier"].nunique())
col2.metric("Countries", filtered_df["Country"].nunique())
col3.metric("Total Incidents", int(filtered_df["Incidents"].sum() if filtered_df["Incidents"].notnull().any() else 0))

# --------- RISK & PERFORMANCE OVERVIEW ---------
st.subheader("Risk and Performance Overview")
col4, col5 = st.columns(2)
with col4:
    st.bar_chart(filtered_df.groupby("Supplier")["Incidents"].sum())
with col5:
    st.bar_chart(filtered_df.groupby("Country")["On-Time Delivery (%)"].mean())

# --------- SUPPLIER DETAILS ---------
st.subheader("Supplier Details")
suppliers = filtered_df["Supplier"].unique()
if len(suppliers) > 0:
    selected_supplier = st.selectbox("Select a supplier for details", suppliers)
    supplier_details = filtered_df[filtered_df["Supplier"] == selected_supplier]
    st.write(supplier_details)

    # --------- RECOMMANDATIONS CLAIRES ET SPECIFIQUES ---------
    recs = []
    high_incidents = supplier_details["Incidents"].max() >= 5
    low_stock = supplier_details["Stock Days"].min() <= 10
    low_delivery = supplier_details["On-Time Delivery (%)"].mean() < 85
    high_risk = supplier_details["Risk Level"].str.contains("High").any()

    if high_incidents:
        recs.append("📍 Incident rate élevé : effectuer un audit qualité, renforcer le suivi et envisager une diversification des fournisseurs.")
    if low_stock:
        recs.append("🔎 Stock de sécurité faible : augmenter le niveau de stock pour éviter les ruptures d'approvisionnement.")
    if low_delivery:
        recs.append("⏰ Taux de livraison à l'heure insuffisant : revoir les processus logistiques et négocier des pénalités.")
    if high_risk:
        recs.append("🚩 Niveau de risque global élevé : planifier des visites sur site et mettre à jour le plan de continuité d'activité.")
    if not recs:
        recs.append("✅ Aucun risque majeur détecté. Maintenir la surveillance et poursuivre les bonnes pratiques.")

    st.markdown("**AI Recommendations:**")
    for r in recs:
        st.write(r)
    nowstr = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    st.markdown(f"_État généré le : **{nowstr}**_")

    # --------- EXPORT POWERPOINT ---------
    import io
    if Presentation is not None:
        ppt_buffer = io.BytesIO()
        prs = Presentation()
        # Slide 1 : Titre
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Rapport fournisseur : {selected_supplier}"
        slide.placeholders[1].text = f"État généré le {nowstr}"

        # Slide 2 : Infos principales
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Données clés du fournisseur"
        text = ""
        for col in supplier_details.columns:
            text += f"- {col} : {', '.join(str(x) for x in supplier_details[col].unique())}\n"
        slide.placeholders[1].text = text

        # Slide 3 : Recommandations
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "AI Recommendations"
        slide.placeholders[1].text = "\n".join(recs)

        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        st.download_button(
            label="Download PowerPoint file",
            data=ppt_buffer,
            file_name=f"{selected_supplier}_report_{nowstr.replace(':','-').replace(' ','_')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    else:
        st.warning(
            "Le module python-pptx n'est pas installé. "
            "Ajoutez `python-pptx` à votre requirements.txt puis redémarrez l'application pour activer l'export PowerPoint."
        )

else:
    st.info("No supplier selected.")

# --------- FOOTER ---------
st.markdown("""
---
Developed for Airbus Commercial Aircraft - by Hynde EL HOUJJAJI - Subcontractor of Capgemini Engineering· [2025]
""")
