import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# --- DATA INITIALIZATION ---

# Sample data for suppliers
SUPPLIER_DATA = {
    'Supplier': ['Supplier A', 'Supplier B', 'Supplier C'],
    'Country': ['France', 'Germany', 'USA'],
    'Risk Score': [0.2, 0.5, 0.8],
    'Risk Level': ['Low', 'Medium', 'High'],
    'Average Delay': [5, 10, 15],
    'Stock Level': [100, 50, 20],
    'Incidents': [1, 3, 5],
    'latitude': [48.8566, 52.5200, 37.7749],
    'longitude': [2.3522, 13.4050, -122.4194]
}
CRITICAL_COMPONENTS = {
    'Supplier A': ['Moteur', 'Avionique'],
    'Supplier B': ['Fuselage', "Train d'atterrissage"],
    'Supplier C': ['Systèmes électriques', 'Hydraulique']
}
PROGRAM_DATA = {
    'Programme': ['A320', 'A350', 'A220'],
    'Risk Score': [0.4, 0.6, 0.3],
    'Average Delay': [7, 12, 5],
    'Stock Level': [80, 40, 60]
}

# --- HELPER FUNCTIONS ---

def compute_risk_level(score: float) -> str:
    if score < 0.33:
        return "Low"
    elif score < 0.66:
        return "Medium"
    else:
        return "High"

def get_supplier_details(df, supplier):
    row = df[df['Supplier'] == supplier].iloc[0]
    return row

def update_df_for_scenario(df: pd.DataFrame, scenario: str) -> pd.DataFrame:
    df = df.copy()
    if scenario == 'Embargo':
        df['Risk Score'] += 0.1
        df['Average Delay'] += 5
    elif scenario == 'Grève':
        df['Risk Score'] += 0.2
        df['Average Delay'] += 10
    elif scenario == 'Guerre':
        df['Risk Score'] += 0.3
        df['Average Delay'] += 15
    # Clamp risk score to 1.0 max
    df['Risk Score'] = df['Risk Score'].clip(upper=1.0)
    df['Risk Level'] = df['Risk Score'].apply(compute_risk_level)
    return df

def export_pptx(supplier_data, logo_path='airbus_logo.png'):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Rapport de gestion des risques fournisseurs"

    # Logo Airbus (optional)
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), height=Inches(1.0))
    else:
        st.warning("Le logo Airbus n'a pas été trouvé et ne sera pas ajouté au rapport.")

    # Détails du fournisseur & recommandations
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(8.5)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = [
        f"Fournisseur : {supplier_data['Supplier']}",
        f"Pays : {supplier_data['Country']}",
        f"Score de risque : {supplier_data['Risk Score']:.2f}",
        f"Niveau de risque : {supplier_data['Risk Level']}",
        f"Délai moyen : {supplier_data['Average Delay']} jours",
        f"Niveau de stock : {supplier_data['Stock Level']}",
        f"Incidents : {supplier_data['Incidents']}",
        "",
        "Recommandations :",
    ]
    if supplier_data['Risk Level'] == 'High':
        lines += ["- Diversifier les fournisseurs", "- Augmenter les stocks de sécurité"]
    elif supplier_data['Risk Level'] == 'Medium':
        lines += ["- Surveiller les performances", "- Planifier des audits réguliers"]
    else:
        lines += ["- Maintenir la collaboration"]

    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18 if line == lines[0] else 14)
        if "Recommandations" in line:
            p.font.bold = True
        if line.startswith("-"):
            p.level = 1

    prs.save('rapport_risques_fournisseurs.pptx')
    st.success('Rapport PowerPoint généré avec succès !')

# --- MAIN APP ---

# Sidebar
st.sidebar.title('Paramètres & Simulation')
scenario = st.sidebar.selectbox(
    'Simulation de scénarios géopolitiques',
    ['Aucun', 'Embargo', 'Grève', 'Guerre']
)
impact = st.sidebar.slider('Impact des actualités sur le risque:', 0.0, 1.0, 0.1)
selected_supplier = st.sidebar.selectbox('Composants critiques par fournisseur', SUPPLIER_DATA['Supplier'])
st.sidebar.write(
    f"**Composants critiques pour {selected_supplier} :**\n"
    f"{', '.join(CRITICAL_COMPONENTS[selected_supplier])}"
)

# Data preparation
df = pd.DataFrame(SUPPLIER_DATA)
df = update_df_for_scenario(df, scenario)
df['Risk Score'] += impact
df['Risk Score'] = df['Risk Score'].clip(upper=1.0)
df['Risk Level'] = df['Risk Score'].apply(compute_risk_level)
df['Predicted Stock Out'] = (df['Average Delay'] * 2).round(1)

# Page Title
st.title('Plateforme de gestion des risques fournisseurs')

# Map with supplier locations and color by risk
colors = df['Risk Level'].map({'Low': '#3CB371', 'Medium': '#FFA500', 'High': '#FF6347'})
st.map(df[['latitude', 'longitude']])

# Supplier selector and details
supplier = st.selectbox('Sélectionnez un fournisseur :', df['Supplier'])
supplier_data = get_supplier_details(df, supplier)
st.markdown(f"""
- **Pays :** {supplier_data['Country']}
- **Score de risque :** {supplier_data['Risk Score']:.2f}
- **Niveau de risque :** {supplier_data['Risk Level']}
- **Délai moyen :** {supplier_data['Average Delay']} jours
- **Niveau de stock :** {supplier_data['Stock Level']}
- **Incidents :** {supplier_data['Incidents']}
- **Predicted Stock Out :** {supplier_data['Predicted Stock Out']} unités
""")

# Main dashboards
st.subheader("Comparatif des fournisseurs")
col1, col2, col3 = st.columns(3)
with col1:
    st.bar_chart(df.set_index('Supplier')['Risk Score'])
with col2:
    st.line_chart(df.set_index('Supplier')['Average Delay'])
with col3:
    st.area_chart(df.set_index('Supplier')['Stock Level'])

# Dashboard par programme avion
st.subheader('Dashboard par programme avion')
df_program = pd.DataFrame(PROGRAM_DATA)
pcol1, pcol2, pcol3 = st.columns(3)
with pcol1:
    st.bar_chart(df_program.set_index('Programme')['Risk Score'])
with pcol2:
    st.line_chart(df_program.set_index('Programme')['Average Delay'])
with pcol3:
    st.area_chart(df_program.set_index('Programme')['Stock Level'])

# Export PowerPoint
if st.button('Exporter un rapport PowerPoint sur le fournisseur sélectionné'):
    export_pptx(supplier_data)

